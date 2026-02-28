#!/usr/bin/env python3
"""
CivicLemma - PowerPoint Presentation Generator
Generates a comprehensive PPTX presentation for the CivicLemma civic engagement platform.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Alias for compatibility
RgbColor = RGBColor

# Color scheme - Professional Blue/Green theme
PRIMARY_COLOR = RgbColor(0x22, 0x8B, 0x22)  # Forest Green
SECONDARY_COLOR = RgbColor(0x1E, 0x90, 0xFF)  # Dodger Blue
ACCENT_COLOR = RgbColor(0xFF, 0x6B, 0x35)  # Orange accent
DARK_COLOR = RgbColor(0x2C, 0x3E, 0x50)  # Dark blue-gray
WHITE_COLOR = RgbColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RgbColor(0xF5, 0xF5, 0xF5)


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide with gradient-like effect"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add colored rectangle as background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()

    # Add title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.33), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE_COLOR
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, title, icon_text=""):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Left colored panel
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = SECONDARY_COLOR
    left_panel.line.fill.background()

    # Icon circle
    if icon_text:
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.5), Inches(2.75), Inches(1.5), Inches(1.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = WHITE_COLOR
        circle.line.fill.background()

        icon_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.1), Inches(1.5), Inches(0.8)
        )
        tf = icon_box.text_frame
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER

    # Title on right
    title_box = slide.shapes.add_textbox(Inches(5), Inches(3), Inches(7.5), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_COLOR

    return slide


def add_content_slide(prs, title, bullet_points, two_columns=False):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    if two_columns and len(bullet_points) > 4:
        # Split into two columns
        mid = len(bullet_points) // 2
        left_points = bullet_points[:mid]
        right_points = bullet_points[mid:]

        # Left column
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.5)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(left_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_COLOR
            p.space_after = Pt(12)

        # Right column
        right_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.5)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(right_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_COLOR
            p.space_after = Pt(12)
    else:
        # Single column
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.6), Inches(12.33), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(22)
            p.font.color.rgb = DARK_COLOR
            p.space_after = Pt(14)

    return slide


def add_architecture_slide(prs):
    """Add architecture diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "System Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Architecture boxes
    components = [
        ("Frontend\n(Next.js 16)", Inches(0.5), Inches(1.8), SECONDARY_COLOR),
        (
            "Backend API\n(Express.js)",
            Inches(4.8),
            Inches(1.8),
            RgbColor(0x9B, 0x59, 0xB6),
        ),
        ("ML Service\n(FastAPI)", Inches(9.1), Inches(1.8), RgbColor(0xE7, 0x4C, 0x3C)),
        (
            "Agent Service\n(FastAPI)",
            Inches(9.1),
            Inches(3.8),
            RgbColor(0xF3, 0x9C, 0x12),
        ),
        (
            "Amazon\nDynamoDB & Cognito",
            Inches(4.8),
            Inches(5.2),
            RgbColor(0xFF, 0xCA, 0x28),
        ),
        (
            "Amazon S3\nImage Storage",
            Inches(0.5),
            Inches(5.2),
            RgbColor(0x00, 0xBC, 0xD4),
        ),
        (
            "Telegram Bot\nIntegration",
            Inches(0.5),
            Inches(3.5),
            RgbColor(0x26, 0xA6, 0x9A),
        ),
    ]

    for text, left, top, color in components:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(1.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(
            left, top + Inches(0.25), Inches(3.8), Inches(0.8)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR
        p.alignment = PP_ALIGN.CENTER

    # Add arrows (simplified with lines)
    # Arrow from Frontend to Backend
    arrow1 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(4.35), Inches(2.2), Inches(0.4), Inches(0.3)
    )
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = DARK_COLOR

    # Arrow from Backend to ML
    arrow2 = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(8.65), Inches(2.2), Inches(0.4), Inches(0.3)
    )
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = DARK_COLOR

    return slide


def add_data_flow_slide(prs):
    """Add data flow diagram slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Data Flow - Issue Reporting"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Flow steps
    steps = [
        ("1. User uploads\nphoto", SECONDARY_COLOR),
        ("2. Image stored\nin S3", RgbColor(0x00, 0xBC, 0xD4)),
        ("3. ML classifies\nissue type", RgbColor(0xE7, 0x4C, 0x3C)),
        ("4. Bedrock generates\ndescription", RgbColor(0x9B, 0x59, 0xB6)),
        ("5. Issue stored\nin DynamoDB", RgbColor(0xFF, 0xCA, 0x28)),
        ("6. Routed to\nMunicipality", PRIMARY_COLOR),
    ]

    start_x = 0.3
    for i, (text, color) in enumerate(steps):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(start_x + i * 2.1),
            Inches(2.5),
            Inches(1.9),
            Inches(1.4),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(
            Inches(start_x + i * 2.1), Inches(2.8), Inches(1.9), Inches(0.9)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR
        p.alignment = PP_ALIGN.CENTER

        # Add arrow between steps
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(start_x + i * 2.1 + 1.95),
                Inches(3.05),
                Inches(0.15),
                Inches(0.25),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_COLOR

    # Bottom note
    note_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(12.33), Inches(2.5)
    )
    tf = note_box.text_frame
    tf.word_wrap = True

    notes = [
        "• Location auto-detected via browser geolocation API",
        "• MobileNetV2 model trained on 9 civic issue categories",
        "• Priority scoring based on issue type, location density, and age",
        "• Real-time updates via DynamoDB streams",
    ]

    for i, note in enumerate(notes):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = note
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(8)

    return slide


def add_tech_stack_slide(prs):
    """Add tech stack slide with categories"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Technology Stack"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Tech categories
    categories = [
        (
            "Frontend",
            ["Next.js 16", "React 18", "TailwindCSS 4", "Radix UI", "TypeScript"],
            SECONDARY_COLOR,
            Inches(0.4),
            Inches(1.5),
        ),
        (
            "Backend",
            [
                "Express.js",
                "TypeScript",
                "AWS SDK",
                "Zod Validation",
                "JWT Auth",
            ],
            RgbColor(0x9B, 0x59, 0xB6),
            Inches(4.6),
            Inches(1.5),
        ),
        (
            "ML Service",
            ["FastAPI", "TensorFlow", "MobileNetV2", "AWS Bedrock", "PIL"],
            RgbColor(0xE7, 0x4C, 0x3C),
            Inches(8.8),
            Inches(1.5),
        ),
        (
            "AI Agent",
            ["AWS Bedrock", "Claude", "Transcribe STT", "Polly TTS", "LangChain"],
            RgbColor(0xF3, 0x9C, 0x12),
            Inches(0.4),
            Inches(4.5),
        ),
        (
            "Infrastructure",
            [
                "Cognito Auth",
                "DynamoDB",
                "S3 + CloudFront",
                "Google Maps",
                "Telegram API",
            ],
            RgbColor(0x00, 0xBC, 0xD4),
            Inches(4.6),
            Inches(4.5),
        ),
        (
            "DevOps",
            ["Vercel", "Render", "Git/GitHub", "npm", "uvicorn"],
            RgbColor(0x2E, 0xCC, 0x71),
            Inches(8.8),
            Inches(4.5),
        ),
    ]

    for cat_name, items, color, left, top in categories:
        # Category box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4), Inches(2.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Category title
        cat_title = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.1), Inches(3.8), Inches(0.4)
        )
        tf = cat_title.text_frame
        p = tf.paragraphs[0]
        p.text = cat_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR

        # Items
        items_box = slide.shapes.add_textbox(
            left + Inches(0.2), top + Inches(0.55), Inches(3.6), Inches(2.1)
        )
        tf = items_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE_COLOR
            p.space_after = Pt(4)

    return slide


def add_issue_types_slide(prs):
    """Add issue types slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Issue Classification Categories (ML Model)"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Issue types grid
    issue_types = [
        ("🚗", "Potholes & Road Damage", "Road surface defects"),
        ("🗑️", "Littering/Garbage", "Public area pollution"),
        ("🅿️", "Illegal Parking", "Parking violations"),
        ("⚠️", "Broken Road Signs", "Damaged signage"),
        ("🌳", "Fallen Trees", "Fallen tree hazards"),
        ("🎨", "Vandalism/Graffiti", "Property defacement"),
        ("🐕", "Dead Animal Pollution", "Animal carcasses"),
        ("🧱", "Damaged Concrete", "Structural damage"),
        ("⚡", "Electrical Damage", "Pole/wire issues"),
    ]

    colors = [
        RgbColor(0xE7, 0x4C, 0x3C),
        RgbColor(0x2E, 0xCC, 0x71),
        RgbColor(0x34, 0x98, 0xDB),
        RgbColor(0xF3, 0x9C, 0x12),
        RgbColor(0x1A, 0xBC, 0x9C),
        RgbColor(0x9B, 0x59, 0xB6),
        RgbColor(0xE9, 0x1E, 0x63),
        RgbColor(0x00, 0xBC, 0xD4),
        RgbColor(0xFF, 0x57, 0x22),
    ]

    for i, (icon, name, desc) in enumerate(issue_types):
        row = i // 3
        col = i % 3
        left = Inches(0.5 + col * 4.2)
        top = Inches(1.5 + row * 1.9)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.9), Inches(1.6)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i]
        box.line.fill.background()

        # Icon and name
        text_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.2), Inches(3.6), Inches(0.5)
        )
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{icon}  {name}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR

        # Description
        desc_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.8), Inches(3.6), Inches(0.6)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE_COLOR

    return slide


def add_user_roles_slide(prs):
    """Add user roles slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "User Roles & Permissions"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    roles = [
        (
            "👤 Citizen (Default)",
            [
                "Report civic issues with photos",
                "Track status of submitted reports",
                "View issues on interactive map",
                "View municipality leaderboard",
                "Access AI voice/chat agent",
            ],
            SECONDARY_COLOR,
            Inches(0.4),
        ),
        (
            "🏛️ Municipality User",
            [
                "View assigned municipality issues",
                "Respond to and update issue status",
                "Upload resolution evidence",
                "Close resolved issues",
                "Access municipality analytics dashboard",
            ],
            RgbColor(0x9B, 0x59, 0xB6),
            Inches(4.6),
        ),
        (
            "⚙️ Platform Admin",
            [
                "Full access to all features",
                "Manage municipalities (CRUD)",
                "View platform-wide statistics",
                "Monitor all issues & responses",
                "User management capabilities",
            ],
            RgbColor(0xE7, 0x4C, 0x3C),
            Inches(8.8),
        ),
    ]

    for title, permissions, color, left in roles:
        # Role box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.5), Inches(4), Inches(5.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Role title
        role_title = slide.shapes.add_textbox(
            left + Inches(0.15), Inches(1.7), Inches(3.7), Inches(0.5)
        )
        tf = role_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR

        # Permissions
        perm_box = slide.shapes.add_textbox(
            left + Inches(0.2), Inches(2.4), Inches(3.6), Inches(4.4)
        )
        tf = perm_box.text_frame
        tf.word_wrap = True
        for i, perm in enumerate(permissions):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"✓ {perm}"
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE_COLOR
            p.space_after = Pt(10)

    return slide


def add_dashboards_slide(prs):
    """Add dashboards overview slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Dashboard Features"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    dashboards = [
        (
            "🗺️ Public Map View",
            [
                "Interactive Google Maps integration",
                "Filter issues by type and status",
                "Cluster view for high-density areas",
                "Real-time issue markers",
            ],
            SECONDARY_COLOR,
            Inches(0.4),
            Inches(1.5),
        ),
        (
            "🏛️ Municipality Dashboard",
            [
                "Issue queue management",
                "Response workflow interface",
                "Resolution image upload",
                "Performance analytics",
                "Monthly trend charts",
            ],
            RgbColor(0x9B, 0x59, 0xB6),
            Inches(6.9),
            Inches(1.5),
        ),
        (
            "📊 Admin Dashboard",
            [
                "Platform-wide statistics",
                "Municipality management",
                "User analytics",
                "Issue category breakdown",
                "System health monitoring",
            ],
            RgbColor(0xE7, 0x4C, 0x3C),
            Inches(0.4),
            Inches(4.3),
        ),
        (
            "🏆 Leaderboard",
            [
                "Municipality rankings",
                "Performance score algorithm",
                "Resolution time tracking",
                "Public accountability",
                "Gamification elements",
            ],
            RgbColor(0xF3, 0x9C, 0x12),
            Inches(6.9),
            Inches(4.3),
        ),
    ]

    for title, features, color, left, top in dashboards:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.1), Inches(2.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Title
        dash_title = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.1), Inches(5.8), Inches(0.4)
        )
        tf = dash_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR

        # Features (2 columns)
        for i, feat in enumerate(features):
            col = i % 2
            row = i // 2
            feat_box = slide.shapes.add_textbox(
                left + Inches(0.1 + col * 3),
                top + Inches(0.55 + row * 0.4),
                Inches(2.9),
                Inches(0.35),
            )
            tf = feat_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {feat}"
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE_COLOR

    return slide


def add_ai_features_slide(prs):
    """Add AI features slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "AI & ML Features"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    ai_features = [
        (
            "🤖 Image Classification",
            "MobileNetV2-based CNN trained on 9 civic issue categories with high accuracy",
            RgbColor(0xE7, 0x4C, 0x3C),
        ),
        (
            "✨ AI Description Generation",
            "AWS Bedrock generates detailed, context-aware issue descriptions",
            RgbColor(0x9B, 0x59, 0xB6),
        ),
        (
            "🎙️ Voice Agent",
            "AWS Bedrock powered conversational agent with Transcribe STT & Polly TTS",
            SECONDARY_COLOR,
        ),
        (
            "💬 Chat Agent",
            "Text-based AI assistant for issue reporting, tracking, and information queries",
            RgbColor(0x2E, 0xCC, 0x71),
        ),
        (
            "📊 Priority Scoring",
            "Intelligent prioritization based on issue type, density, location, and age factors",
            RgbColor(0xF3, 0x9C, 0x12),
        ),
        (
            "🤖 Telegram Bot",
            "Full-featured bot integration for reporting issues via Telegram messaging app",
            RgbColor(0x00, 0xBC, 0xD4),
        ),
    ]

    for i, (title, desc, color) in enumerate(ai_features):
        row = i // 2
        col = i % 2
        left = Inches(0.4 + col * 6.5)
        top = Inches(1.5 + row * 1.9)

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.2), Inches(1.7)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Title
        feat_title = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.15), Inches(5.9), Inches(0.4)
        )
        tf = feat_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR

        # Description
        desc_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.65), Inches(5.9), Inches(0.9)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE_COLOR

    return slide


def add_api_slide(prs):
    """Add API endpoints slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RESTful API Endpoints"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    # Table-like structure
    endpoints = [
        (
            "PUBLIC ENDPOINTS",
            [
                "GET  /api/issues              - List issues with filtering",
                "POST /api/issues              - Create new issue report",
                "GET  /api/issues/:id          - Get specific issue",
                "GET  /api/municipalities      - List municipalities",
                "GET  /api/leaderboard         - Municipality rankings",
            ],
            SECONDARY_COLOR,
        ),
        (
            "PROTECTED ENDPOINTS",
            [
                "POST /api/issues/:id/respond  - Municipality response",
                "POST /api/issues/:id/close    - Close resolved issue",
                "POST /api/upload/presign      - S3 presigned URL",
                "GET  /api/admin/stats         - Platform statistics",
            ],
            RgbColor(0xE7, 0x4C, 0x3C),
        ),
        (
            "ML & AGENT ENDPOINTS",
            [
                "POST /classify                - Classify issue image",
                "POST /generate-description    - AI description",
                "POST /agent/chat              - Chat with AI agent",
                "POST /agent/voice             - Voice interaction",
            ],
            RgbColor(0x9B, 0x59, 0xB6),
        ),
    ]

    top = Inches(1.4)
    for cat_name, items, color in endpoints:
        # Category header
        cat_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), top, Inches(12.5), Inches(0.45)
        )
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = color
        cat_box.line.fill.background()

        cat_title = slide.shapes.add_textbox(
            Inches(0.5), top + Inches(0.08), Inches(12.3), Inches(0.35)
        )
        tf = cat_title.text_frame
        p = tf.paragraphs[0]
        p.text = cat_name
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR

        top += Inches(0.5)

        # Endpoints
        for item in items:
            item_box = slide.shapes.add_textbox(
                Inches(0.5), top, Inches(12.3), Inches(0.35)
            )
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(13)
            p.font.name = "Consolas"
            p.font.color.rgb = DARK_COLOR
            top += Inches(0.35)

        top += Inches(0.15)

    return slide


def add_security_slide(prs):
    """Add security features slide"""
    return add_content_slide(
        prs,
        "Security & Authentication",
        [
            "Cognito Authentication with JWT token verification",
            "Role-based access control (RBAC) for all endpoints",
            "Secure file upload via S3 presigned URLs",
            "Input validation using Zod schemas on all requests",
            "CORS protection with whitelisted origins",
            "Environment-based configuration (no hardcoded secrets)",
            "AWS IAM policies for server-side operations",
            "Secure session management with token refresh",
        ],
    )


def add_deployment_slide(prs):
    """Add deployment architecture slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY_COLOR
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Deployment Architecture"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR

    deployments = [
        (
            "🌐 Frontend\nVercel",
            "Next.js SSR, Edge Functions,\nAutomatic CI/CD",
            SECONDARY_COLOR,
            Inches(0.5),
            Inches(1.8),
        ),
        (
            "⚙️ Backend\nRender",
            "Express.js, Auto-scaling,\nDocker containers",
            RgbColor(0x9B, 0x59, 0xB6),
            Inches(4.6),
            Inches(1.8),
        ),
        (
            "🧠 ML Service\nAWS",
            "FastAPI, TensorFlow,\nPython 3.9+",
            RgbColor(0xE7, 0x4C, 0x3C),
            Inches(8.7),
            Inches(1.8),
        ),
        (
            "🤖 Agent Service\nAWS",
            "FastAPI, AWS Bedrock,\nTelegram Bot",
            RgbColor(0xF3, 0x9C, 0x12),
            Inches(0.5),
            Inches(4.5),
        ),
        (
            "☁️ Amazon",
            "DynamoDB,\nCognito Auth",
            RgbColor(0xFF, 0xCA, 0x28),
            Inches(4.6),
            Inches(4.5),
        ),
        (
            "📦 Amazon S3",
            "Image Storage,\nCloudFront CDN",
            RgbColor(0x00, 0xBC, 0xD4),
            Inches(8.7),
            Inches(4.5),
        ),
    ]

    for title, desc, color, left, top in deployments:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(2.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Title
        dep_title = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(0.15), Inches(3.6), Inches(0.8)
        )
        tf = dep_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(
            left + Inches(0.1), top + Inches(1.1), Inches(3.6), Inches(1)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE_COLOR
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_future_scope_slide(prs):
    """Add future scope slide"""
    return add_content_slide(
        prs,
        "Future Enhancements",
        [
            "Multi-language support for broader accessibility",
            "Push notifications for issue status updates",
            "Augmented Reality (AR) for issue visualization",
            "Predictive analytics for infrastructure maintenance",
            "Integration with government e-governance portals",
            "Citizen reward/incentive system for reporting",
            "Advanced analytics dashboard with ML insights",
            "Mobile app (React Native) for offline-first experience",
        ],
        two_columns=True,
    )


def add_thank_you_slide(prs):
    """Add thank you slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_COLOR
    bg.line.fill.background()

    # Thank you text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(12.33), Inches(1)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE_COLOR
    p.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5), Inches(12.33), Inches(1)
    )
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "CivicLemma - Empowering Citizens, Transforming Governance"
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE_COLOR
    p.alignment = PP_ALIGN.CENTER

    return slide


def create_presentation():
    """Create the complete CivicLemma presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs, "CivicLemma", "AI-Powered Civic Engagement Platform for India")

    # Slide 2: Problem Statement
    add_section_slide(prs, "Problem Statement", "❓")

    # Slide 3: Problem Details
    add_content_slide(
        prs,
        "The Challenge",
        [
            "Citizens lack efficient channels to report local infrastructure issues",
            "Manual reporting processes are time-consuming and often ignored",
            "No transparency in issue resolution by municipalities",
            "Lack of accountability metrics for local government bodies",
            "Difficulty in categorizing and prioritizing civic issues",
            "No centralized system for tracking issue lifecycle",
            "Communication gap between citizens and municipal authorities",
        ],
    )

    # Slide 4: Solution Overview
    add_section_slide(prs, "Our Solution", "💡")

    # Slide 5: Solution Features
    add_content_slide(
        prs,
        "CivicLemma - Key Features",
        [
            "📷 Photo-based issue reporting with automatic location detection",
            "🤖 AI-powered image classification using MobileNetV2 CNN",
            "✨ Intelligent issue descriptions via AWS Bedrock AI",
            "🗺️ Interactive map view for browsing all reported issues",
            "🏛️ Dedicated municipality dashboard for issue management",
            "🏆 Public leaderboard ranking municipalities by performance",
            "🎙️ Voice & chat AI agent for conversational reporting",
            "📱 Telegram bot integration for mobile-first users",
        ],
    )

    # Slide 6: Tech Stack
    add_tech_stack_slide(prs)

    # Slide 7: Architecture
    add_architecture_slide(prs)

    # Slide 8: Data Flow
    add_data_flow_slide(prs)

    # Slide 9: Issue Types
    add_issue_types_slide(prs)

    # Slide 10: User Roles
    add_user_roles_slide(prs)

    # Slide 11: Dashboards
    add_dashboards_slide(prs)

    # Slide 12: AI Features
    add_ai_features_slide(prs)

    # Slide 13: API Endpoints
    add_api_slide(prs)

    # Slide 14: Security
    add_security_slide(prs)

    # Slide 15: Deployment
    add_deployment_slide(prs)

    # Slide 16: Future Scope
    add_future_scope_slide(prs)

    # Slide 17: Thank You
    add_thank_you_slide(prs)

    return prs


def main():
    """Main function to generate the presentation"""
    print("🎨 Generating CivicLemma PowerPoint Presentation...")

    prs = create_presentation()

    output_path = "CivicLemma_Presentation.pptx"
    prs.save(output_path)

    print(f"✅ Presentation saved to: {os.path.abspath(output_path)}")
    print(f"📊 Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
